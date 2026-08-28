"""
Projektideen — Portfolio-Folie (PPTX) und Cloud-Illustration (OpenAI Images).

PPTX nutzt lokale Bewertungsdaten (ai_*). Cloud-Illustration: nur DSGVO-geprüfter
Prompt an OpenAI — kein Rohtext, keine Personennamen.
"""
from __future__ import annotations

import html as html_lib
import io
import json
import logging
import os
import re
import unicodedata
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .m01_config import get_settings
from .m03_db import get_session
from .m08_llm import get_available_models, have_key, try_models_with_messages, model_supports_vision
from .m16_idea import ProjectIdea, get_idea

log = logging.getLogger(__name__)

OPENAI_IMAGE_MODELS: dict[str, str] = {
    "dall-e-3": "dall-e-3",
    "dall-e-2": "dall-e-2",
    "gpt-image-1": "gpt-image-1",
    "gpt-image-1-mini": "gpt-image-1-mini",
}
DEFAULT_OPENAI_IMAGE_MODEL = "dall-e-3"

# Provider für Visualisierung (Prompt/Folienstruktur) — optional getrennt von KI-Einstellungen
VISUAL_TEXT_PROVIDERS: tuple[str, ...] = ("anthropic", "ollama", "openai")
VISUAL_TEXT_MODELS: dict[str, list[str]] = {
    "openai": ["gpt-5.4", "gpt-5.4-mini", "gpt-4o", "gpt-4o-mini"],
    "anthropic": ["sonnet-4.6", "opus-4.6", "haiku-4.5"],
    "ollama": [
        "qwen2.5vl:7b",
        "qwen2.5vl:32b",
        "llava:13b",
        "llava:34b",
        "qwen3:32b",
        "llama3.3:70b",
        "qwen3:8b",
        "llama3.2",
    ],
}
VISUAL_TEXT_DEFAULT_MODELS: dict[str, str] = {
    "openai": "gpt-5.4",
    "anthropic": "sonnet-4.6",
    "ollama": "qwen3:32b",
}

IDEA_VISUAL_OUTPUT_FORMATS: dict[str, str] = {
    "html": "HTML-Bericht — interaktiv (neuer Tab)",
    "pptx": "PowerPoint — Download + Prozessdiagramm",
    "png_local": "PNG — Prozessdiagramm / Canvas (lokal)",
    "png_cloud": "PNG — Cloud-Illustration (OpenAI Images)",
    "docx": "Word — Bericht nur Text",
    "docx_png": "Word — Bericht + Diagramm",
}

_JSON_BLOCK_RE = re.compile(r"\{[\s\S]*\}")

CLOUD_LLM_PROVIDERS = frozenset({"openai", "anthropic"})

_DIAGRAM_FONT_CANDIDATES = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "C:/Windows/Fonts/arial.ttf",
    "/Library/Fonts/Arial.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
]
_diagram_font_path_cache: Optional[str] = None
_diagram_font_warned = False


def _nfc_text(text: str) -> str:
    return unicodedata.normalize("NFC", text) if text else ""


def _resolve_diagram_font_path() -> Optional[str]:
    global _diagram_font_path_cache
    if _diagram_font_path_cache is not None:
        return _diagram_font_path_cache or None
    env_path = os.environ.get("SLIT_DIAGRAM_FONT", "").strip()
    candidates: list[str] = []
    if env_path:
        candidates.append(env_path)
    candidates.extend(_DIAGRAM_FONT_CANDIDATES)
    candidates.append("arial.ttf")
    for path in candidates:
        if path and Path(path).is_file():
            _diagram_font_path_cache = path
            return path
    _diagram_font_path_cache = ""
    return None


def _load_diagram_font(size: int) -> ImageFont.ImageFont:
    global _diagram_font_warned
    path = _resolve_diagram_font_path()
    if path:
        try:
            return ImageFont.truetype(path, size)
        except OSError as exc:
            log.warning("Diagram font %s failed: %s", path, exc)
    if not _diagram_font_warned:
        log.warning("No TrueType diagram font found; umlauts may render incorrectly")
        _diagram_font_warned = True
    return ImageFont.load_default()


def _draw_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.ImageFont,
    **kwargs: Any,
) -> None:
    draw.text(xy, _nfc_text(text), font=font, **kwargs)


@dataclass
class DeckContent:
    title: str
    subtitle: str = ""
    summary_lines: list[str] = field(default_factory=list)
    resource_lines: list[str] = field(default_factory=list)
    challenge_lines: list[str] = field(default_factory=list)
    phase_lines: list[str] = field(default_factory=list)
    phase_details: list[dict[str, Any]] = field(default_factory=list)
    recommendation_lines: list[str] = field(default_factory=list)


def visual_text_providers_available() -> list[str]:
    return [p for p in VISUAL_TEXT_PROVIDERS if have_key(p)]


def visual_text_models_map() -> dict[str, list[str]]:
    result: dict[str, list[str]] = {}
    for p in visual_text_providers_available():
        curated = list(VISUAL_TEXT_MODELS.get(p, []))
        if p == "ollama":
            live = get_available_models("ollama")
            models = [m for m in curated if m in live]
            if not models and live:
                models = live[:12]
            elif models:
                for m in live:
                    if m not in models and len(models) < 12:
                        models.append(m)
        else:
            models = curated
        if models:
            result[p] = models
    return result


def visual_vision_models_map() -> dict[str, list[str]]:
    """Modelle mit Vision-Unterstützung (für Referenz-Uploads)."""
    result: dict[str, list[str]] = {}
    for p, models in visual_text_models_map().items():
        vision = [m for m in models if model_supports_vision(p, m)]
        if vision:
            result[p] = vision
    return result


def resolve_visual_llm(
    visual_provider: str | None,
    visual_model: str | None,
    fallback_provider: str,
    fallback_model: str,
) -> tuple[str, str]:
    vp = (visual_provider or "").strip().lower()
    if vp and vp in VISUAL_TEXT_PROVIDERS and have_key(vp):
        models = VISUAL_TEXT_MODELS.get(vp, [])
        vm = (visual_model or "").strip()
        if vm and vm in models:
            return vp, vm
        default = VISUAL_TEXT_DEFAULT_MODELS.get(vp, "")
        if default and default in models:
            return vp, default
        return vp, models[0] if models else fallback_model
    return fallback_provider, fallback_model or ""


def is_cloud_llm_provider(provider: str) -> bool:
    return (provider or "").strip().lower() in CLOUD_LLM_PROVIDERS


def default_local_assess_llm() -> tuple[str, str]:
    """Bevorzugt Ollama-Textmodell für Bewertung (kein Cloud-Risiko)."""
    if not have_key("ollama"):
        return "", ""
    live = get_available_models("ollama")
    preferred = VISUAL_TEXT_DEFAULT_MODELS.get("ollama", "")
    if preferred and preferred in live and not model_supports_vision("ollama", preferred):
        return "ollama", preferred
    for cand in live:
        if not model_supports_vision("ollama", cand):
            return "ollama", cand
    for cand in live:
        return "ollama", cand
    return "ollama", ""


def idea_assess_provider_defaults(settings: dict[str, Any]) -> tuple[str, str]:
    lp, lm = default_local_assess_llm()
    if lp:
        return lp, lm
    return (settings.get("provider") or "openai"), (settings.get("model") or "")


def validate_assess_cloud_gates(
    idea: ProjectIdea,
    assess_provider: str,
    assess_model: str,
    input_provider: str,
    source_tasks: set[str],
    cloud_confirm: bool,
    vision_cloud_confirm: bool,
) -> Optional[str]:
    return validate_cloud_gates_for_references(
        assess_provider,
        assess_model,
        input_provider,
        source_tasks,
        cloud_confirm,
        vision_cloud_confirm,
        idea=idea,
    )


def _cloud_gate_attachment_state(
    idea: ProjectIdea | None = None,
    ref_bundle: Any | None = None,
) -> tuple[bool, bool]:
    from .m17_visual_lab_refs import LabReferenceBundle

    has_att = False
    has_images = False
    bundle: LabReferenceBundle | None = None
    if idea is not None:
        from .m16_idea import _idea_source_bundle, _load_stored_attachments

        has_att = bool(_load_stored_attachments(idea)) or bool(
            (idea.source_reference_text or "").strip()
        )
        bundle = _idea_source_bundle(idea)
    if ref_bundle is not None:
        has_att = has_att or bool(ref_bundle.stored or ref_bundle.text_blocks)
        if bundle is None or ref_bundle.images:
            bundle = ref_bundle
    if bundle is not None:
        has_images = bool(bundle.image_payload())
    return has_att, has_images


def validate_cloud_gates_for_references(
    assess_provider: str,
    assess_model: str,
    input_provider: str,
    source_tasks: set[str],
    cloud_confirm: bool,
    vision_cloud_confirm: bool,
    idea: ProjectIdea | None = None,
    ref_bundle: Any | None = None,
) -> Optional[str]:
    from .m08_llm import get_model_id

    has_att, has_images = _cloud_gate_attachment_state(idea, ref_bundle)

    assess_cloud = is_cloud_llm_provider(assess_provider)
    input_cloud = is_cloud_llm_provider(input_provider)
    uses_cloud = assess_cloud or (
        input_cloud
        and (
            "vision_describe" in source_tasks
            or ("vision_images" in source_tasks and has_images)
        )
    )

    if has_att and uses_cloud and not cloud_confirm:
        return "cloud_confirm"

    vision_cloud_needed = False
    if has_images and "vision_images" in source_tasks and assess_cloud:
        mid = get_model_id(assess_provider, assess_model) or assess_model
        if model_supports_vision(assess_provider, mid):
            vision_cloud_needed = True
    if has_images and "vision_describe" in source_tasks and input_cloud:
        vision_cloud_needed = True

    if vision_cloud_needed and not vision_cloud_confirm:
        return "vision_cloud_confirm"
    return None


def _merge_refinement_with_reference(
    refinement_notes: str,
    reference_text: str,
    for_cloud: bool = False,
) -> str:
    parts: list[str] = []
    if (reference_text or "").strip():
        ref = reference_text.strip()[:8000]
        if for_cloud:
            ref = sanitize_for_cloud_text(ref)
        parts.append("Berücksichtige folgende Referenzunterlagen:\n" + ref)
    if (refinement_notes or "").strip():
        note = refinement_notes.strip()
        if for_cloud:
            note = sanitize_for_cloud_text(note)
        parts.append(note)
    return "\n\n".join(parts)


def resolve_idea_reference_context(
    idea: ProjectIdea,
    source_tasks: set[str],
    input_provider: str,
    input_model: str,
    for_cloud: bool = False,
) -> tuple[str, list[tuple[bytes, str]]]:
    from .m16_idea import _idea_source_bundle
    from .m17_visual_lab_refs import (
        DEFAULT_SOURCE_TASKS,
        describe_reference_images,
        filter_bundle_for_source_tasks,
    )

    tasks = source_tasks if source_tasks is not None else set(DEFAULT_SOURCE_TASKS)
    bundle = _idea_source_bundle(idea)
    if not bundle:
        return "", []
    filtered = filter_bundle_for_source_tasks(bundle, tasks)
    text_parts: list[str] = []
    if "extract_text" in tasks:
        t = filtered.merged_text().strip()
        if not t and idea.source_reference_text:
            t = (idea.source_reference_text or "").strip()
        if t:
            text_parts.append(sanitize_for_cloud_text(t) if for_cloud else t)
    if "vision_describe" in tasks and filtered.images:
        desc = describe_reference_images(filtered, input_provider, input_model)
        if desc:
            raw = sanitize_for_cloud_text(desc) if for_cloud else desc
            text_parts.append("Referenz-Bildbeschreibung:\n" + raw)
    images = filtered.image_payload() if "vision_images" in tasks else []
    return "\n\n".join(text_parts), images


def sanitize_structured_field(text: str) -> str:
    """Kontakt und Herr/Frau — ohne Paar-Heuristik (deutsche Produktnamen bleiben erhalten)."""
    from .m20_pii_stage1 import sanitize_structured_field as _stage1_structured

    return _stage1_structured(text)


def sanitize_for_cloud_text(text: str) -> str:
    """Stufe 1: Regex-Heuristiken. Stufe 2: swiss-pii-anonymizer (falls installiert)."""
    sanitized, _ = sanitize_for_cloud_with_meta(text)
    return sanitized


def sanitize_for_cloud_with_meta(text: str) -> tuple[str, list[dict[str, str | float]]]:
    """Wie sanitize_for_cloud_text, zusätzlich Presidio-Findings aus einem anonymize-Lauf."""
    if not text:
        return "", []
    from .m18_cloud_pii import apply_swiss_pii_anonymize_details
    from .m20_pii_stage1 import apply_pii_stage1

    t = apply_pii_stage1(text, preserve_newlines=False)
    return apply_swiss_pii_anonymize_details(t)


def idea_images_dir() -> Path:
    d = Path(get_settings().data_dir) / "idea_images"
    d.mkdir(parents=True, exist_ok=True)
    return d


def idea_decks_dir() -> Path:
    d = Path(get_settings().data_dir) / "idea_decks"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _now() -> datetime:
    return datetime.now(timezone.utc)


def _parse_json_dict(raw: str | None) -> dict[str, Any]:
    if not raw:
        return {}
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text, count=1)
        text = re.sub(r"\s*```\s*$", "", text)
    m = _JSON_BLOCK_RE.search(text)
    if not m:
        return {}
    try:
        data = json.loads(m.group(0))
        return data if isinstance(data, dict) else {}
    except json.JSONDecodeError:
        return {}


_ILLUSTRATION_SYSTEM = (
    "Du erstellst einen ENGLISCHEN Bild-Prompt für DALL-E (abstrakte Portfolio-Illustration).\n"
    "DSGVO / Datenschutz — STRIKT:\n"
    "- KEINE Personennamen, keine Anreden mit Namen, keine E-Mail, Telefon, Adressen.\n"
    "- KEINE konkreten internen Vorgänge, Aktenzeichen, IDs oder Rohtext der Idee.\n"
    "- Erlaubt: generischer Projekttyp, abstrakte Symbole, Farben, Stimmung.\n"
    "- Firmen-/Behörden-/Fachbereichsname DARF einmal vorkommen wenn sachlich nützlich.\n"
    "- Keine fotorealistischen Personen, keine Gesichter.\n"
    "Stil: professionell, öffentliche Verwaltung Schweiz, minimal, blau-grau.\n"
    "Antwort: NUR der Prompt-Text (max. 900 Zeichen), ohne Anführungszeichen oder Markdown."
)


def _build_illustration_user_prompt(idea: ProjectIdea, refinement_notes: str = "") -> str:
    parts: list[str] = []
    if idea.ai_project_name:
        parts.append(f"Projektname (generisch): {sanitize_structured_field(idea.ai_project_name)}")
    if idea.fachabteilung:
        parts.append(f"Fachbereich/Organisation: {sanitize_structured_field(idea.fachabteilung)}")
    if idea.ai_summary:
        parts.append(f"Sachliche Kurzfassung: {sanitize_for_cloud_text(idea.ai_summary)}")
    if idea.ai_recommendation:
        parts.append(f"Empfehlung (abstrakt): {sanitize_for_cloud_text(idea.ai_recommendation)}")
    if not parts:
        parts.append("Generisches öffentliches Verwaltungsprojekt, Innovation, abstrakt.")
    ref = sanitize_for_cloud_text(refinement_notes)
    if ref:
        parts.append(f"Anpassung für neue Illustration (ohne Personennamen): {ref}")
    return "\n".join(parts)


def build_dsgvo_illustration_prompt(
    idea: ProjectIdea,
    llm_provider: str = "openai",
    llm_model: str = "",
    refinement_notes: str = "",
    extra_reference: str = "",
) -> Optional[str]:
    user = _build_illustration_user_prompt(idea, refinement_notes)
    if (extra_reference or "").strip():
        user += "\n\nReferenzunterlagen (abstrakt):\n" + sanitize_for_cloud_text(extra_reference)[:3000]
    raw = try_models_with_messages(
        llm_provider,
        _ILLUSTRATION_SYSTEM,
        [{"role": "user", "content": user}],
        max_tokens=500,
        temperature=0.4,
        model=llm_model or None,
    )
    prompt = sanitize_for_cloud_text((raw or "").strip().strip('"').strip("'"))
    if len(prompt) < 20:
        title = sanitize_structured_field(idea.ai_project_name or idea.title or "Public sector project")
        dept = sanitize_structured_field(idea.fachabteilung or "")
        ref = sanitize_for_cloud_text(refinement_notes)
        prompt = (
            f"Abstract minimalist illustration for a Swiss public administration portfolio concept: "
            f"{title}. "
            f"{f'Context: {dept}. ' if dept else ''}"
            f"{f'Adjustment: {ref}. ' if ref else ''}"
            "Professional blue and grey tones, geometric icons, no people, no text, no logos."
        )
    return prompt[:900]


def generate_openai_illustration(prompt: str, image_model: str = DEFAULT_OPENAI_IMAGE_MODEL) -> Optional[bytes]:
    if not have_key("openai"):
        log.warning("OpenAI key missing for illustration")
        return None
    model_id = OPENAI_IMAGE_MODELS.get(image_model, image_model)
    try:
        from openai import OpenAI

        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))
    except Exception as exc:
        log.warning("OpenAI client init failed: %s", exc)
        return None

    try:
        if model_id == "dall-e-3":
            resp = client.images.generate(
                model=model_id,
                prompt=prompt,
                size="1792x1024",
                quality="standard",
                n=1,
            )
        elif model_id == "dall-e-2":
            resp = client.images.generate(model=model_id, prompt=prompt, size="1024x1024", n=1)
        else:
            resp = client.images.generate(model=model_id, prompt=prompt, size="1024x1024", n=1)
        item = resp.data[0]
        if getattr(item, "b64_json", None):
            import base64

            return base64.b64decode(item.b64_json)
        if getattr(item, "url", None):
            import urllib.request

            with urllib.request.urlopen(item.url, timeout=120) as r:
                return r.read()
    except Exception as exc:
        log.warning("OpenAI image generation failed: %s", exc)
        return None
    return None


def deck_content_from_idea(
    idea: ProjectIdea,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
) -> DeckContent:
    base = DeckContent(
        title=idea.ai_project_name or idea.title or f"Projektidee #{idea.id}",
        subtitle=" · ".join(
            [x for x in [
                idea.fachabteilung,
                "Projektportfolio — Ersteinschätzung",
                idea.created_at.strftime("%d.%m.%Y"),
            ] if x]
        ),
        summary_lines=[s.strip() for s in re.split(r"(?<=[.!?])\s+", idea.ai_summary or "") if s.strip()][:6],
        resource_lines=_resource_lines_from_idea(idea),
        challenge_lines=[
            f"{c.get('title', '')}: {c.get('description', '')} [{c.get('severity', '')}]"
            for c in idea.challenges[:5]
        ],
        phase_lines=[
            f"{p.get('name', '')} ({p.get('duration_estimate', '')}): {p.get('description', '')}"
            for p in idea.phases[:6]
        ],
        phase_details=[
            {
                "title": _clean_phase_title(p.get("name") or f"Phase {i+1}")[:80],
                "bullets": [str(p.get("description") or "")[:200]] if p.get("description") else [],
                "parallel_note": str(p.get("duration_estimate") or ""),
            }
            for i, p in enumerate(idea.phases[:6])
        ],
        recommendation_lines=[idea.ai_recommendation] if idea.ai_recommendation else [],
    )
    ref = (refinement_notes or "").strip()
    if not ref:
        return base

    system = (
        "Du passt Portfolio-Folieninhalt an. Antwort NUR als JSON:\n"
        '{"title":"...","subtitle":"...","summary_lines":[],"resource_lines":[],'
        '"challenge_lines":[],"phase_lines":[],"recommendation_lines":[]}\n'
        "Wende die Anpassungswünsche an, bleibe sachlich. Keine neuen Personennamen erfinden."
    )
    user = (
        f"Bisheriger Inhalt:\n{json.dumps(base.__dict__, ensure_ascii=False)}\n\n"
        f"Anpassungswünsche:\n{ref}"
    )
    raw = try_models_with_messages(
        llm_provider,
        system,
        [{"role": "user", "content": user}],
        max_tokens=1200,
        temperature=0.3,
        model=llm_model or None,
    )
    parsed = _parse_json_dict(raw)
    if not parsed:
        base.recommendation_lines = base.recommendation_lines + [f"Anpassung: {ref}"]
        return base
    return DeckContent(
        title=(parsed.get("title") or base.title)[:120],
        subtitle=(parsed.get("subtitle") or base.subtitle)[:200],
        summary_lines=_as_str_list(parsed.get("summary_lines")) or base.summary_lines,
        resource_lines=_as_str_list(parsed.get("resource_lines")) or base.resource_lines,
        challenge_lines=_as_str_list(parsed.get("challenge_lines")) or base.challenge_lines,
        phase_lines=_as_str_list(parsed.get("phase_lines")) or base.phase_lines,
        phase_details=_deck_content_from_parsed(parsed).phase_details or base.phase_details,
        recommendation_lines=_as_str_list(parsed.get("recommendation_lines")) or base.recommendation_lines,
    )


def _as_str_list(val: Any) -> list[str]:
    if not isinstance(val, list):
        return []
    return [str(x).strip() for x in val if str(x).strip()]


def _resource_lines_from_idea(idea: ProjectIdea) -> list[str]:
    lines: list[str] = []
    if idea.ai_internal_pt is not None:
        lines.append(f"Intern: {idea.ai_internal_pt:.0f} Personentage")
        if idea.ai_internal_pt_reasoning:
            lines.append(idea.ai_internal_pt_reasoning)
    if idea.ai_external_cost is not None:
        lines.append(f"Extern: CHF {idea.ai_external_cost:,.0f}".replace(",", "'"))
        if idea.ai_external_cost_reasoning:
            lines.append(idea.ai_external_cost_reasoning)
    return lines


_LAB_DECK_SYSTEM = (
    "Erstelle reichhaltigen Portfolio-/Schulungs-Folieninhalt (öffentliche Verwaltung Schweiz).\n"
    "Antwort NUR als JSON:\n"
    '{"title":"kurzer Titel","subtitle":"optional","summary_lines":["max 2 Sätze"],'
    '"resource_lines":[],"challenge_lines":[],"phase_lines":["1. Name: Kurz"],'
    '"phase_details":[{"title":"Initiierung","bullets":["Ziel klären","Stakeholder","Go/No-Go"],'
    '"parallel_note":""}],"recommendation_lines":[]}\n'
    "Regeln:\n"
    "- phase_details: je Phase 2–4 konkrete Stichpunkte (wie Schulungsfolie).\n"
    "- phase_lines: nummerierte Kurzzeile pro Phase (für Diagramm).\n"
    "- Keine Personennamen. Sachlich, visuell strukturierbar.\n"
    "- summary_lines NICHT Phasen wiederholen."
)


def _phase_details_from_lines(phase_lines: list[str]) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for line in phase_lines:
        raw = line.strip()
        if raw.startswith("•"):
            raw = raw[1:].strip()
        if raw and raw[0].isdigit() and "." in raw[:4]:
            raw = raw.split(".", 1)[1].strip()
        title = raw
        bullets: list[str] = []
        if ":" in raw:
            title, rest = raw.split(":", 1)
            title = title.strip()
            parts = [p.strip() for p in re.split(r"[·;]", rest) if p.strip()]
            bullets = parts if parts else [rest.strip()]
        out.append({"title": _clean_phase_title(title)[:80], "bullets": bullets[:6], "parallel_note": ""})
    return out


def _deck_content_from_parsed(parsed: dict[str, Any], fallback_title: str = "Visual") -> DeckContent:
    phase_lines = _as_str_list(parsed.get("phase_lines"))
    raw_details = parsed.get("phase_details")
    phase_details: list[dict[str, Any]] = []
    if isinstance(raw_details, list):
        for item in raw_details[:8]:
            if not isinstance(item, dict):
                continue
            title = _clean_phase_title(str(item.get("title") or ""))
            if not title:
                continue
            bullets = _as_str_list(item.get("bullets"))
            parallel = str(item.get("parallel_note") or item.get("parallel") or "").strip()
            phase_details.append({"title": title[:80], "bullets": bullets[:6], "parallel_note": parallel[:160]})
    if not phase_details and phase_lines:
        phase_details = _phase_details_from_lines(phase_lines)
    return DeckContent(
        title=(parsed.get("title") or fallback_title)[:120],
        subtitle=(parsed.get("subtitle") or "")[:200],
        summary_lines=_as_str_list(parsed.get("summary_lines")),
        resource_lines=_as_str_list(parsed.get("resource_lines")),
        challenge_lines=_as_str_list(parsed.get("challenge_lines")),
        phase_lines=phase_lines,
        phase_details=phase_details,
        recommendation_lines=_as_str_list(parsed.get("recommendation_lines")),
    )


def _fallback_deck_from_prompt(prompt: str) -> DeckContent:
    lines = [ln.strip() for ln in (prompt or "").splitlines() if ln.strip()]
    title = lines[0][:80] if lines else "Visual"
    body = lines[1:] if len(lines) > 1 else [prompt[:400]]
    phase_lines = [f"{i + 1}. {ln[:100]}" for i, ln in enumerate(body[:8])]
    return DeckContent(
        title=title,
        subtitle="SlitProjektHub Visual-Lab",
        phase_lines=phase_lines,
        phase_details=_phase_details_from_lines(phase_lines),
    )


def _phase_labels_from_lines(phase_lines: list[str]) -> list[str]:
    labels: list[str] = []
    for line in phase_lines[:8]:
        raw = line.strip()
        if raw.startswith("•"):
            raw = raw[1:].strip()
        if raw and raw[0].isdigit() and "." in raw[:4]:
            raw = raw.split(".", 1)[1].strip()
        label = raw.split(":")[0].split("(")[0].strip()
        labels.append(label[:48] or raw[:48])
    return labels


def build_process_diagram_png(labels: list[str], title: str = "") -> bytes:
    """Horizontale Prozessdarstellung (Phasen/Steps) als PNG."""
    w, h = 1280, 720
    img = Image.new("RGB", (w, h), color=(245, 247, 252))
    draw = ImageDraw.Draw(img)
    font_t = _load_diagram_font(36)
    font_b = _load_diagram_font(18)
    font_s = _load_diagram_font(14)
    title = _nfc_text(title)
    if title:
        _draw_text(draw, (48, 36), title[:70], font_t, fill=(30, 58, 95))
    n = max(len(labels), 1)
    box_w, box_h, gap = 200, 88, 36
    total = n * box_w + (n - 1) * gap
    start_x = max(48, (w - total) // 2)
    y = 200
    colors = [(30, 58, 95), (45, 85, 130), (60, 110, 165), (75, 130, 190)]
    for i, label in enumerate(labels[:8]):
        label = _nfc_text(label)
        x = start_x + i * (box_w + gap)
        col = colors[i % len(colors)]
        draw.rounded_rectangle([x, y, x + box_w, y + box_h], radius=12, fill=col)
        num = str(i + 1)
        _draw_text(draw, (x + 14, y + 10), num, font_b, fill=(255, 255, 255))
        words = label.split()
        line, ly = "", y + 38
        for word in words:
            test = f"{line} {word}".strip()
            if len(test) > 22 and line:
                _draw_text(draw, (x + 12, ly), line, font_s, fill=(230, 235, 245))
                ly += 18
                line = word
            else:
                line = test
        if line:
            _draw_text(draw, (x + 12, ly), line[:26], font_s, fill=(230, 235, 245))
        if i < len(labels) - 1:
            ax = x + box_w + 6
            draw.polygon([(ax, y + box_h // 2 - 8), (ax + gap - 12, y + box_h // 2), (ax, y + box_h // 2 + 8)], fill=(120, 140, 170))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


_PHASE_COLORS = [
    (232, 222, 248),
    (214, 245, 228),
    (214, 232, 248),
    (252, 228, 210),
    (248, 214, 226),
    (230, 240, 250),
]
_PHASE_BORDER = [
    (109, 76, 180),
    (46, 125, 80),
    (37, 99, 168),
    (196, 120, 58),
    (180, 72, 110),
    (90, 120, 150),
]

_PHASE_BOX_GAP = 28


_LEADING_PHASE_RE = re.compile(
    r"^(?:phase\s+)?\d+\s*[\.\)\:\-–—•·\u2022\u00b7\uf0b7]\s*",
    re.IGNORECASE,
)


def _clean_phase_title(title: str) -> str:
    """Entfernt führende Nummerierung ('1. …', 'Phase 2: …'), die sonst doppelt erscheint."""
    t = _nfc_text((title or "").strip())
    t = re.sub(r"^phase\s+\d+\s*:\s*", "", t, flags=re.IGNORECASE).strip()
    t = _LEADING_PHASE_RE.sub("", t).strip()
    return t[:80]


def _measure_text_width(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> int:
    bbox = draw.textbbox((0, 0), _nfc_text(text), font=font)
    return max(0, bbox[2] - bbox[0])


def _wrap_to_width(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: ImageFont.ImageFont,
    max_width: int,
) -> list[str]:
    text = _nfc_text(text or "").strip()
    if not text:
        return []
    words = text.split()
    lines: list[str] = []
    cur = ""
    for word in words:
        test = f"{cur} {word}".strip()
        if not cur or _measure_text_width(draw, test, font) <= max_width:
            cur = test
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines or [text]


def _vertical_layout(
    details: list[dict[str, Any]],
    title: str,
    w: int = 1200,
) -> tuple[list[dict[str, Any]], int]:
    """Berechnet umbrochene Zeilen und Boxhöhen, bevor das PNG gerendert wird."""
    dummy = Image.new("RGB", (8, 8), color=(255, 255, 255))
    draw = ImageDraw.Draw(dummy)
    font_h = _load_diagram_font(22)
    font_b = _load_diagram_font(16)
    font_s = _load_diagram_font(13)
    inner = w - 80 * 2 - 48
    laid: list[dict[str, Any]] = []
    for i, pd in enumerate(details[:8]):
        heading = f"Phase {i + 1}: {_clean_phase_title(str(pd.get('title') or ''))}"
        heading_lines = _wrap_to_width(draw, heading, font_h, inner)
        bullet_lines: list[str] = []
        for b in (pd.get("bullets") or [])[:6]:
            wrapped = _wrap_to_width(draw, f"· {b}", font_b, inner)
            bullet_lines.extend(wrapped[:4])
        par = (pd.get("parallel_note") or "").strip()
        par_lines = _wrap_to_width(draw, f"parallel: {par}", font_s, inner) if par else []
        box_h = 20 + len(heading_lines) * 26 + len(bullet_lines) * 22 + len(par_lines) * 18 + 16
        laid.append({
            "height": max(72, box_h),
            "heading_lines": heading_lines,
            "bullet_lines": bullet_lines,
            "par_lines": par_lines,
            "color_i": i,
        })
    top = 80 if title else 40
    body = sum(x["height"] for x in laid)
    gaps = _PHASE_BOX_GAP * max(0, len(laid) - 1)
    return laid, max(400, top + body + gaps + 48)


def build_vertical_process_diagram_png(
    phase_details: list[dict[str, Any]],
    title: str = "",
) -> bytes:
    w = 1200
    font_t = _load_diagram_font(32)
    font_h = _load_diagram_font(22)
    font_b = _load_diagram_font(16)
    font_s = _load_diagram_font(13)
    title = _nfc_text(title)
    laid, h = _vertical_layout(phase_details, title, w)
    img = Image.new("RGB", (w, h), color=(252, 252, 254))
    draw = ImageDraw.Draw(img)
    if title:
        _draw_text(draw, (40, 24), title[:90], font_t, fill=(30, 58, 95))
    y = 80 if title else 40
    cx = w // 2
    for i, box in enumerate(laid):
        ph = box["height"]
        col = _PHASE_COLORS[box["color_i"] % len(_PHASE_COLORS)]
        border = _PHASE_BORDER[box["color_i"] % len(_PHASE_BORDER)]
        bx = 80
        draw.rounded_rectangle([bx, y, w - 80, y + ph], radius=14, fill=col, outline=border, width=2)
        ty = y + 12
        for line in box["heading_lines"]:
            _draw_text(draw, (bx + 20, ty), line, font_h, fill=border)
            ty += 26
        ty += 4
        for line in box["bullet_lines"]:
            _draw_text(draw, (bx + 22, ty), line, font_b, fill=(55, 65, 80))
            ty += 22
        for line in box["par_lines"]:
            _draw_text(draw, (bx + 22, ty), line, font_s, fill=(100, 110, 130))
            ty += 18
        if i < len(laid) - 1:
            ay = y + ph + 4
            draw.polygon([(cx, ay + 18), (cx - 10, ay), (cx + 10, ay)], fill=(140, 150, 170))
            y += ph + _PHASE_BOX_GAP
        else:
            y += ph
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _render_slide_panel(heading: str, lines: list[str], width: int = 1280) -> Image.Image:
    font_h = _load_diagram_font(26)
    font_b = _load_diagram_font(15)
    lh = 26
    h = 70 + min(len(lines), 12) * lh + 30
    img = Image.new("RGB", (width, h), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, width, 4], fill=(30, 58, 95))
    _draw_text(draw, (48, 28), heading, font_h, fill=(30, 58, 95))
    y = 68
    for line in lines[:12]:
        _draw_text(draw, (56, y), f"• {line[:110]}", font_b, fill=(50, 55, 65))
        y += lh
    return img


def _render_title_panel(content: DeckContent, width: int = 1280) -> Image.Image:
    font_l = _load_diagram_font(40)
    font_s = _load_diagram_font(20)
    img = Image.new("RGB", (width, 220), color=(30, 58, 95))
    draw = ImageDraw.Draw(img)
    _draw_text(draw, (48, 60), content.title[:80], font_l, fill=(255, 255, 255))
    if content.subtitle:
        _draw_text(draw, (48, 130), content.subtitle[:120], font_s, fill=(200, 210, 225))
    return img


def build_deck_composite_preview_png(content: DeckContent) -> bytes:
    panels: list[Image.Image] = []
    panels.append(_render_title_panel(content))
    if content.summary_lines:
        panels.append(_render_slide_panel("Zusammenfassung", content.summary_lines))
    details = content.phase_details or _phase_details_from_lines(content.phase_lines)
    if len(details) >= 2:
        vimg = Image.open(io.BytesIO(build_vertical_process_diagram_png(details, content.title[:60])))
        panels.append(vimg)
    if content.phase_lines:
        panels.append(_render_slide_panel("Grobe Phasenplanung", content.phase_lines))
    for pd in details[:4]:
        bullets = pd.get("bullets") or []
        if bullets:
            panels.append(_render_slide_panel(f"Detail: {pd.get('title', '')}", bullets))
    if content.recommendation_lines:
        panels.append(_render_slide_panel("Empfehlung", content.recommendation_lines))
    panels.append(_render_slide_panel("Hinweis", ["KI-Vorbewertung — ersetzt keine fachliche Prüfung · SlitProjektHub"]))
    gap = 12
    width = 1280
    total_h = sum(p.height for p in panels) + gap * max(0, len(panels) - 1)
    out = Image.new("RGB", (width, total_h), color=(236, 239, 244))
    y = 0
    for p in panels:
        if p.width != width:
            resample = getattr(Image, "Resampling", Image).LANCZOS
            p = p.resize((width, int(p.height * width / p.width)), resample)
        out.paste(p, (0, y))
        y += p.height + gap
    buf = io.BytesIO()
    out.save(buf, format="PNG")
    return buf.getvalue()


def deck_content_from_lab_prompt(
    prompt: str,
    llm_provider: str = "openai",
    llm_model: str = "",
    reference_text: str = "",
    reference_images: list[tuple[bytes, str]] | None = None,
) -> Optional[DeckContent]:
    user_content = prompt.strip()
    if reference_text.strip():
        user_content += "\n\n--- Referenzmaterial (extrahiert) ---\n" + reference_text.strip()[:10000]
    images = reference_images or []
    model_id = llm_model or ""
    use_images = images and model_supports_vision(llm_provider, model_id)
    raw = try_models_with_messages(
        llm_provider,
        _LAB_DECK_SYSTEM,
        [{"role": "user", "content": user_content}],
        max_tokens=2000,
        temperature=0.45,
        model=llm_model or None,
        images=images if use_images else None,
    )
    if not raw:
        return _fallback_deck_from_prompt(user_content)
    parsed = _parse_json_dict(raw)
    if not parsed:
        return _fallback_deck_from_prompt(user_content)
    return _deck_content_from_parsed(parsed, fallback_title=user_content.strip().split("\n")[0][:80] or "Visual Test")


def _add_title_slide(prs: Presentation, content: DeckContent) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(9), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = content.title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    if content.subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(9), Inches(1))
        sp = sub.text_frame.paragraphs[0]
        sp.text = content.subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_diagram_slide(prs: Presentation, png_bytes: bytes, heading: str = "Prozessdarstellung") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
    hp = head.text_frame.paragraphs[0]
    hp.text = heading
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.35), Inches(1.0), width=Inches(9.3))


def _add_bullet_slide(prs: Presentation, heading: str, lines: list[str]) -> None:
    if not lines:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.6))
    hp = head.text_frame.paragraphs[0]
    hp.text = heading
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(9), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines[:12]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line.startswith("•") else f"• {line}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)


def _add_phase_detail_slide(prs: Presentation, phase: dict[str, Any], index: int) -> None:
    bullets = phase.get("bullets") or []
    if not bullets:
        return
    title = f"Phase {index + 1}: {_clean_phase_title(str(phase.get('title', '')))}"
    lines = list(bullets)
    par = (phase.get("parallel_note") or "").strip()
    if par:
        lines.append(f"(Parallel: {par})")
    _add_bullet_slide(prs, title[:80], lines)


def build_pptx_bytes(content: DeckContent) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs, content)
    if content.summary_lines:
        _add_bullet_slide(prs, "Zusammenfassung", content.summary_lines)
    if content.resource_lines:
        _add_bullet_slide(prs, "Ressourcen & Kosten", content.resource_lines)
    if content.challenge_lines:
        _add_bullet_slide(prs, "Herausforderungen", content.challenge_lines)
    if content.phase_lines:
        details = content.phase_details or _phase_details_from_lines(content.phase_lines)
        labels = _phase_labels_from_lines(content.phase_lines)
        if len(details) >= 2:
            _add_diagram_slide(
                prs,
                build_vertical_process_diagram_png(details, content.title[:60]),
                "Prozess / Phasen",
            )
        elif len(labels) >= 2:
            _add_diagram_slide(
                prs,
                build_process_diagram_png(labels, content.title[:60]),
                "Prozess / Phasen",
            )
        for i, pd in enumerate(details[:5]):
            _add_phase_detail_slide(prs, pd, i)
        _add_bullet_slide(prs, "Grobe Phasenplanung", content.phase_lines)
    if content.recommendation_lines:
        _add_bullet_slide(prs, "Empfehlung", content.recommendation_lines)
    foot = prs.slides.add_slide(prs.slide_layouts[6])
    fp = foot.shapes.add_textbox(Inches(0.6), Inches(3), Inches(9), Inches(1)).text_frame.paragraphs[0]
    fp.text = "KI-Vorbewertung — ersetzt keine fachliche Prüfung · SlitProjektHub"
    fp.font.size = Pt(11)
    fp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_deck_preview_png(content: DeckContent) -> bytes:
    return build_deck_composite_preview_png(content)


def idea_html_dir() -> Path:
    d = Path(get_settings().data_dir) / "idea_html"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _html_esc(text: str) -> str:
    return html_lib.escape(_nfc_text(text or ""), quote=True)


def _html_list(items: list[str]) -> str:
    if not items:
        return "<p class=\"muted\">—</p>"
    lis = "".join(f"<li>{_html_esc(x)}</li>" for x in items)
    return f"<ul>{lis}</ul>"


def build_html_report(content: DeckContent) -> str:
    """Selbstständige HTML-Datei: Inhaltsverzeichnis, Sprungmarken, volle Textbreite."""
    details = content.phase_details or _phase_details_from_lines(content.phase_lines)
    nav_items: list[tuple[str, str]] = [("top", content.title or "Bericht")]
    if content.summary_lines:
        nav_items.append(("summary", "Zusammenfassung"))
    if content.resource_lines:
        nav_items.append(("resources", "Ressourcen & Kosten"))
    if content.challenge_lines:
        nav_items.append(("challenges", "Herausforderungen"))
    if details:
        nav_items.append(("phases", "Phasen"))
        for i, pd in enumerate(details):
            nav_items.append((f"phase-{i+1}", f"Phase {i+1}"))
    if content.recommendation_lines:
        nav_items.append(("recommendation", "Empfehlung"))

    nav_html = "".join(
        f'<a href="#{hid}">{_html_esc(label)}</a>' for hid, label in nav_items
    )

    phase_html_parts: list[str] = []
    colors = ["#6d4cb4", "#2e7d50", "#2563a8", "#c4783a", "#b4486e", "#5a7896"]
    bgs = ["#ece4f8", "#d6f5e4", "#d6e8f8", "#fce4d2", "#f8d6e2", "#e6f0fa"]
    for i, pd in enumerate(details):
        col = colors[i % len(colors)]
        bg = bgs[i % len(bgs)]
        bullets = "".join(f"<li>{_html_esc(b)}</li>" for b in (pd.get("bullets") or []) if b)
        par = (pd.get("parallel_note") or "").strip()
        par_html = f'<p class="parallel">Parallel / Dauer: {_html_esc(par)}</p>' if par else ""
        phase_html_parts.append(
            f'<article class="phase" id="phase-{i+1}" style="--accent:{col};--bg:{bg}">'
            f"<h3>Phase {i + 1}: {_html_esc(_clean_phase_title(str(pd.get('title') or '')))}</h3>"
            f"{'<ul>' + bullets + '</ul>' if bullets else ''}"
            f"{par_html}</article>"
        )

    sections: list[str] = []
    if content.summary_lines:
        sections.append(
            f'<section id="summary"><h2>Zusammenfassung</h2>{_html_list(content.summary_lines)}</section>'
        )
    if content.resource_lines:
        sections.append(
            f'<section id="resources"><h2>Ressourcen &amp; Kosten</h2>{_html_list(content.resource_lines)}</section>'
        )
    if content.challenge_lines:
        sections.append(
            f'<section id="challenges"><h2>Herausforderungen</h2>{_html_list(content.challenge_lines)}</section>'
        )
    if phase_html_parts:
        sections.append(
            '<section id="phases"><h2>Phasenplanung</h2>'
            f'<div class="phases">{"".join(phase_html_parts)}</div></section>'
        )
    if content.recommendation_lines:
        sections.append(
            f'<section id="recommendation"><h2>Empfehlung</h2>{_html_list(content.recommendation_lines)}</section>'
        )

    title = _html_esc(content.title or "Projektbericht")
    subtitle = _html_esc(content.subtitle or "")
    return f"""<!doctype html>
<html lang="de">
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>{title}</title>
<style>
:root {{
  --ink:#1e293b; --muted:#64748b; --line:#e2e8f0; --bg:#f8fafc; --card:#fff; --accent:#1e3a5f;
}}
* {{ box-sizing:border-box; }}
html {{ scroll-behavior:smooth; }}
body {{
  margin:0; font-family: ui-sans-serif, system-ui, "Segoe UI", sans-serif;
  color:var(--ink); background:var(--bg); line-height:1.55;
}}
.layout {{ display:grid; grid-template-columns: 240px 1fr; min-height:100vh; }}
nav {{
  position:sticky; top:0; height:100vh; overflow:auto;
  background:var(--accent); color:#e2e8f0; padding:1.4rem 1rem;
}}
nav .brand {{ font-size:.72rem; letter-spacing:.08em; text-transform:uppercase; color:#94a3b8; margin-bottom:.75rem; }}
nav a {{
  display:block; color:#e2e8f0; text-decoration:none; font-size:.88rem;
  padding:.35rem .5rem; border-radius:6px; margin-bottom:.15rem;
}}
nav a:hover, nav a:focus {{ background:rgba(255,255,255,.12); }}
main {{ padding:2rem 2.5rem 4rem; max-width:1100px; }}
header#top {{ margin-bottom:1.75rem; }}
h1 {{ font-size:2rem; line-height:1.2; margin:0 0 .35rem; color:var(--accent); }}
.subtitle {{ color:var(--muted); margin:0; font-size:1rem; }}
h2 {{ font-size:1.25rem; color:var(--accent); margin:2rem 0 .75rem; border-bottom:2px solid var(--line); padding-bottom:.35rem; }}
h3 {{ margin:0 0 .5rem; font-size:1.05rem; }}
ul {{ margin:.35rem 0 0; padding-left:1.2rem; }}
li {{ margin:.25rem 0; }}
.muted {{ color:var(--muted); }}
.phases {{ display:grid; gap:1rem; }}
.phase {{
  background:var(--bg); border-left:5px solid var(--accent); border-radius:10px;
  padding:1rem 1.15rem; box-shadow:0 1px 2px rgba(15,23,42,.06);
}}
.phase .parallel {{ color:var(--muted); font-size:.9rem; margin:.5rem 0 0; }}
footer {{ margin-top:2.5rem; color:var(--muted); font-size:.8rem; border-top:1px solid var(--line); padding-top:.75rem; }}
@media (max-width: 880px) {{
  .layout {{ grid-template-columns: 1fr; }}
  nav {{ position:relative; height:auto; }}
  main {{ padding:1.25rem; }}
}}
@media print {{
  nav {{ display:none; }}
  .layout {{ display:block; }}
  main {{ max-width:none; }}
}}
</style>
</head>
<body>
<div class="layout">
  <nav>
    <div class="brand">SlitProjektHub</div>
    {nav_html}
  </nav>
  <main>
    <header id="top">
      <h1>{title}</h1>
      {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
    </header>
    {''.join(sections)}
    <footer>KI-Vorbewertung — ersetzt keine fachliche Prüfung · SlitProjektHub</footer>
  </main>
</div>
</body>
</html>
"""


def _persist_html_report(idea_id: int, content: DeckContent) -> str:
    name = f"report_{idea_id}_{uuid.uuid4().hex[:10]}.html"
    (idea_html_dir() / name).write_text(build_html_report(content), encoding="utf-8")
    return name


def _unlink_html(obj: ProjectIdea) -> None:
    if not obj.html_path:
        return
    p = idea_html_dir() / Path(obj.html_path).name
    if p.is_file():
        p.unlink(missing_ok=True)


def generate_portfolio_deck(
    idea_id: int,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
) -> Optional[ProjectIdea]:
    idea = get_idea(idea_id)
    if not idea or idea.status != "bewertet":
        return None
    content = deck_content_from_idea(idea, refinement_notes, llm_provider, llm_model)
    deck_name = f"deck_{idea_id}_{uuid.uuid4().hex[:10]}.pptx"
    preview_name = f"deck_preview_{idea_id}.png"
    html_name = _persist_html_report(idea_id, content)
    (idea_decks_dir() / deck_name).write_bytes(build_pptx_bytes(content))
    (idea_images_dir() / preview_name).write_bytes(build_deck_preview_png(content))
    with get_session() as ses:
        obj = ses.get(ProjectIdea, idea_id)
        if not obj:
            return None
        _unlink_html(obj)
        obj.deck_path = deck_name
        obj.deck_preview_path = preview_name
        obj.deck_generated_at = _now()
        obj.html_path = html_name
        obj.html_generated_at = _now()
        obj.updated_at = _now()
        ses.add(obj)
        ses.commit()
        ses.refresh(obj)
        return obj


def generate_cloud_illustration(
    idea_id: int,
    image_model: str = DEFAULT_OPENAI_IMAGE_MODEL,
    llm_provider: str = "openai",
    llm_model: str = "",
    refinement_notes: str = "",
    extra_reference: str = "",
) -> Optional[ProjectIdea]:
    idea = get_idea(idea_id)
    if not idea or idea.status != "bewertet":
        return None
    safe_prompt = build_dsgvo_illustration_prompt(
        idea,
        llm_provider,
        llm_model,
        refinement_notes=refinement_notes,
        extra_reference=extra_reference,
    )
    if not safe_prompt:
        return None
    img_bytes = generate_openai_illustration(safe_prompt, image_model)
    if not img_bytes:
        return None
    fname = f"ill_{idea_id}_{uuid.uuid4().hex[:10]}.png"
    (idea_images_dir() / fname).write_bytes(img_bytes)
    with get_session() as ses:
        obj = ses.get(ProjectIdea, idea_id)
        if not obj:
            return None
        obj.image_path = fname
        obj.image_source = "dalle"
        obj.illustration_model = OPENAI_IMAGE_MODELS.get(image_model, image_model)
        obj.illustration_prompt_safe = safe_prompt[:500]
        obj.illustration_generated_at = _now()
        obj.updated_at = _now()
        ses.add(obj)
        ses.commit()
        ses.refresh(obj)
        return obj


def idea_docx_dir() -> Path:
    d = Path(get_settings().data_dir) / "idea_docx"
    d.mkdir(parents=True, exist_ok=True)
    return d


def build_docx_bytes(content: DeckContent, include_diagram: bool = False) -> bytes:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_heading(content.title, 0)
    if content.subtitle:
        doc.add_paragraph(content.subtitle)

    def _bullets(heading: str, lines: list[str]) -> None:
        if not lines:
            return
        doc.add_heading(heading, level=1)
        for line in lines:
            doc.add_paragraph(line, style="List Bullet")

    _bullets("Zusammenfassung", content.summary_lines)
    _bullets("Ressourcen & Kosten", content.resource_lines)
    _bullets("Herausforderungen", content.challenge_lines)
    _bullets("Grobe Phasenplanung", content.phase_lines)
    _bullets("Empfehlung", content.recommendation_lines)

    if include_diagram and content.phase_lines:
        labels = _phase_labels_from_lines(content.phase_lines)
        if len(labels) >= 2:
            doc.add_heading("Prozessdarstellung", level=1)
            png = build_process_diagram_png(labels, content.title[:60])
            tmp = idea_images_dir() / f"_tmp_{uuid.uuid4().hex[:8]}.png"
            tmp.write_bytes(png)
            try:
                doc.add_picture(str(tmp), width=Inches(6.0))
            finally:
                tmp.unlink(missing_ok=True)

    doc.add_paragraph("KI-Vorbewertung — ersetzt keine fachliche Prüfung · SlitProjektHub")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generate_local_diagram_png(
    idea_id: int,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
) -> Optional[ProjectIdea]:
    idea = get_idea(idea_id)
    if not idea or idea.status != "bewertet":
        return None
    content = deck_content_from_idea(idea, refinement_notes, llm_provider, llm_model)
    fname = f"diag_{idea_id}_{uuid.uuid4().hex[:10]}.png"
    (idea_images_dir() / fname).write_bytes(build_deck_preview_png(content))
    with get_session() as ses:
        obj = ses.get(ProjectIdea, idea_id)
        if not obj:
            return None
        obj.image_path = fname
        obj.image_source = "diagram"
        obj.illustration_model = None
        obj.illustration_prompt_safe = None
        obj.illustration_generated_at = _now()
        obj.updated_at = _now()
        ses.add(obj)
        ses.commit()
        ses.refresh(obj)
        return obj


def generate_idea_docx(
    idea_id: int,
    include_diagram: bool = False,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
) -> Optional[ProjectIdea]:
    idea = get_idea(idea_id)
    if not idea or idea.status != "bewertet":
        return None
    content = deck_content_from_idea(idea, refinement_notes, llm_provider, llm_model)
    docx_name = f"report_{idea_id}_{uuid.uuid4().hex[:10]}.docx"
    (idea_docx_dir() / docx_name).write_bytes(build_docx_bytes(content, include_diagram=include_diagram))
    html_name = _persist_html_report(idea_id, content)
    with get_session() as ses:
        obj = ses.get(ProjectIdea, idea_id)
        if not obj:
            return None
        _unlink_html(obj)
        obj.docx_path = docx_name
        obj.docx_generated_at = _now()
        obj.html_path = html_name
        obj.html_generated_at = _now()
        obj.updated_at = _now()
        ses.add(obj)
        ses.commit()
        ses.refresh(obj)
        return obj


def generate_idea_html(
    idea_id: int,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
) -> Optional[ProjectIdea]:
    idea = get_idea(idea_id)
    if not idea or idea.status != "bewertet":
        return None
    content = deck_content_from_idea(idea, refinement_notes, llm_provider, llm_model)
    html_name = _persist_html_report(idea_id, content)
    preview_name = f"deck_preview_{idea_id}.png"
    (idea_images_dir() / preview_name).write_bytes(build_deck_preview_png(content))
    with get_session() as ses:
        obj = ses.get(ProjectIdea, idea_id)
        if not obj:
            return None
        _unlink_html(obj)
        obj.html_path = html_name
        obj.html_generated_at = _now()
        obj.deck_preview_path = preview_name
        obj.updated_at = _now()
        ses.add(obj)
        ses.commit()
        ses.refresh(obj)
        return obj


def generate_idea_visual(
    idea_id: int,
    output_format: str,
    refinement_notes: str = "",
    llm_provider: str = "openai",
    llm_model: str = "",
    image_model: str = DEFAULT_OPENAI_IMAGE_MODEL,
    input_llm_provider: str = "",
    input_llm_model: str = "",
    source_tasks: set[str] | None = None,
) -> tuple[Optional[ProjectIdea], Optional[str]]:
    fmt = (output_format or "").strip().lower()
    if fmt not in IDEA_VISUAL_OUTPUT_FORMATS:
        return None, "invalid_format"

    idea = get_idea(idea_id)
    if not idea:
        return None, "not_found"
    from .m17_visual_lab_refs import DEFAULT_SOURCE_TASKS

    src = source_tasks if source_tasks is not None else set(DEFAULT_SOURCE_TASKS)
    ip, im = resolve_visual_llm(input_llm_provider, input_llm_model, llm_provider, llm_model)
    out_cloud = is_cloud_llm_provider(llm_provider)
    ref_text, _ = resolve_idea_reference_context(idea, src, ip, im, for_cloud=out_cloud)
    merged_notes = _merge_refinement_with_reference(
        refinement_notes, ref_text, for_cloud=out_cloud,
    )

    if fmt == "html":
        obj = generate_idea_html(idea_id, merged_notes, llm_provider, llm_model)
        return obj, None if obj else "generation_failed"
    if fmt == "pptx":
        obj = generate_portfolio_deck(idea_id, merged_notes, llm_provider, llm_model)
        return obj, None if obj else "generation_failed"
    if fmt == "png_local":
        obj = generate_local_diagram_png(idea_id, merged_notes, llm_provider, llm_model)
        return obj, None if obj else "generation_failed"
    if fmt == "png_cloud":
        if not have_key("openai"):
            return None, "no_key"
        obj = generate_cloud_illustration(
            idea_id,
            image_model,
            llm_provider,
            llm_model,
            refinement_notes,
            extra_reference=ref_text,
        )
        return obj, None if obj else "png_failed"
    if fmt == "docx":
        obj = generate_idea_docx(idea_id, False, merged_notes, llm_provider, llm_model)
        return obj, None if obj else "generation_failed"
    if fmt == "docx_png":
        obj = generate_idea_docx(idea_id, True, merged_notes, llm_provider, llm_model)
        return obj, None if obj else "generation_failed"
    return None, "invalid_format"
